import os
import sys
import importlib.util
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import pandas as pd

# Supported image extensions that go through OCR
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"}


def extract_text(filepath):
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()

        elif ext == ".pdf":
            reader = PdfReader(filepath)
            return " ".join([page.extract_text() or "" for page in reader.pages])

        elif ext == ".docx":
            doc = Document(filepath)
            return " ".join([p.text for p in doc.paragraphs])

        elif ext == ".pptx":
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return " ".join(text)

        elif ext in [".xlsx", ".xls"]:
            df = pd.read_excel(filepath)
            return df.to_string()

        elif ext in IMAGE_EXTENSIONS:
            return _extract_text_from_image(filepath)

        else:
            return ""

    except Exception as e:
        print(f"❌ Error reading {filepath}: {e}")
        return ""


def _load_smart_ocr():
    """
    Load SmartOCR from ocr.py using its absolute path so Django
    can find it regardless of working directory or PYTHONPATH.
    ocr.py must live in the same folder as this file (the core/ app folder).
    """
    # core/ocr.py sits next to this file
    this_dir = os.path.dirname(os.path.abspath(__file__))
    ocr_path = os.path.join(this_dir, "ocr.py")

    if not os.path.isfile(ocr_path):
        raise FileNotFoundError(f"ocr.py not found at: {ocr_path}")

    spec   = importlib.util.spec_from_file_location("ocr", ocr_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module.SmartOCR


def _extract_text_from_image(filepath: str) -> str:
    """Run SmartOCR on an image and return extracted plain text."""
    try:
        SmartOCR = _load_smart_ocr()
        ocr      = SmartOCR()
        result   = ocr.process(filepath)
        return result.full_text.strip()
    except FileNotFoundError as e:
        print(f"❌ {e}")
        return ""
    except Exception as e:
        print(f"❌ OCR failed for {filepath}: {e}")
        return ""