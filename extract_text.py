import os
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text(file_path):
    ext = file_path.lower().split('.')[-1]

    try:
        if ext == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif ext == "pdf":
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text

        elif ext == "docx":
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == "pptx":
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            return ""

    except Exception as e:
        print(f"❌ Error reading {file_path}: {e}")
        return ""